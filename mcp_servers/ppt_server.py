from fastmcp import FastMCP
from pptx import Presentation
from pathlib import Path
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

mcp = FastMCP("ppt-server")

prs = None


def get_theme_styles(theme: str) -> dict:
    themes = {
        "Classic Blue": {
             "bg": RGBColor(15, 32, 64),  
             "title": RGBColor(255, 255, 255), 
             "text": RGBColor(220, 230, 241),  
             "subtitle": RGBColor(180, 200, 220),
             "accent": RGBColor(91, 155, 213),
        },
        "Modern Dark": {
            "bg": RGBColor(34, 34, 34),
            "title": RGBColor(255, 255, 255),
            "text": RGBColor(230, 230, 230),
            "subtitle": RGBColor(200, 200, 200),
            "accent": RGBColor(0, 176, 240),
        },
        "Minimal Light": {
            "bg": RGBColor(248, 249, 250),
            "title": RGBColor(33, 37, 41),
            "text": RGBColor(60, 60, 60),
            "subtitle": RGBColor(100, 100, 100),
            "accent": RGBColor(180, 180, 180),
        },
        "Green Professional": {
            "bg": RGBColor(255, 255, 255),
            "title": RGBColor(0, 97, 0),
            "text": RGBColor(30, 30, 30),
            "subtitle": RGBColor(90, 90, 90),
            "accent": RGBColor(112, 173, 71),
        },
        "Purple Gradient": {
            "bg": RGBColor(245, 240, 255),
            "title": RGBColor(88, 28, 135),
            "text": RGBColor(40, 40, 40),
            "subtitle": RGBColor(100, 80, 120),
            "accent": RGBColor(155, 89, 182),
        },
    }

    return themes.get(theme, themes["Classic Blue"])


def apply_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_accent_bar(slide, prs, accent_color):
    shape = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, 300000
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()


@mcp.tool()
def create_presentation() -> str:
    """
    Create a new PowerPoint presentation.

    Returns:
        Success message confirming a new presentation was created.

    Notes:
        Call this before adding any slides.
    """
    global prs
    prs = Presentation()
    return "Presentation created"


@mcp.tool()
def add_title_slide(title: str, subtitle: str = "", theme: str = "Classic Blue") -> str:
    """
    Add a themed title slide to the presentation.

    Args:
        title: Main title for the presentation.
        subtitle: Subtitle text for the title slide.
        theme: Theme name selected by the user.

    Returns:
        Success message with the slide title.
    """
    global prs

    if prs is None:
        raise ValueError("Presentation not created. Call create_presentation first.")

    styles = get_theme_styles(theme)

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    apply_background(slide, styles["bg"])
    add_top_accent_bar(slide, prs, styles["accent"])

    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.size = Pt(30)
    title_para.font.bold = True
    title_para.font.color.rgb = styles["title"]
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_para = slide.placeholders[1].text_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = styles["subtitle"]
    subtitle_para.alignment = PP_ALIGN.CENTER

    return f"Title slide added: {title}"


@mcp.tool()
def add_bullet_slide(title: str, bullets: list[str], theme: str = "Classic Blue") -> str:
    """
    Add a themed bullet-point content slide.

    Args:
        title: Slide title.
        bullets: List of bullet points for the slide body.
        theme: Theme name selected by the user.

    Returns:
        Success message with the slide title.
    """
    global prs

    if prs is None:
        raise ValueError("Presentation not created. Call create_presentation first.")

    styles = get_theme_styles(theme)

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    apply_background(slide, styles["bg"])
    add_top_accent_bar(slide, prs, styles["accent"])

    slide.shapes.title.text = title

    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = styles["title"]

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = styles["text"]
        p.level = 0

    return f"Slide added: {title}"


@mcp.tool()
def save_presentation(filename: str = "generated_presentation.pptx") -> str:
    """
    Save the current presentation to the output folder.

    Args:
        filename: Name of the output PowerPoint file.

    Returns:
        Full file path of the saved presentation.
    """
    global prs

    if prs is None:
        raise ValueError("Presentation not created. Call create_presentation first.")

    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)

    path = output_dir / filename
    prs.save(path)

    return str(path)


if __name__ == "__main__":
    mcp.run()